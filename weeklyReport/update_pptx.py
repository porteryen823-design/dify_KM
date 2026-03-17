import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
import datetime

def replace_text_in_shape(shape, old_text, new_text):
    # Deep Text Extraction
    if hasattr(shape, "text_frame") and shape.text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)
                    print(f"Replaced text in shape '{shape.name}': {old_text} -> {new_text}")
        
    # Table logic
    if hasattr(shape, "has_table") and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                            print(f"Replaced text in table '{shape.name}': {old_text} -> {new_text}")
                
    # Group recursion
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            replace_text_in_shape(sub_shape, old_text, new_text)

def update_presentation(pptx_path, output_path, replacements):
    prs = pptx.Presentation(pptx_path)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            for old_text, new_text in replacements.items():
                replace_text_in_shape(shape, old_text, new_text)
                
    prs.save(output_path)
    print(f"Successfully saved modified presentation to: {output_path}")

if __name__ == "__main__":
    original_file = r'c:\VSCode_Proj\Dify\weeklyReport\WeekReport_20260316.pptx'
    modified_file = r'c:\VSCode_Proj\Dify\weeklyReport\WeekReport_Modified.pptx'
    
    # 抓取目前系統日期
    today_str = datetime.datetime.now().strftime("%Y.%m.%d")
    
    # 將原檔案中的 "2026.03.16" 和假設的 placeholder 替換為今日
    replacements = {
        "2026.03.16": today_str,
        "###Date1####": today_str
    }
    
    update_presentation(original_file, modified_file, replacements)
